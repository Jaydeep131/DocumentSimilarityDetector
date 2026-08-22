import os
import shutil
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --------------------------------------------------
# CONFIG & FOLDER SETUP
# --------------------------------------------------
project_dir = r"d:\Document\study\Degree\sem 5\PDS\projects\DocumentSimilarityDetector"
assets_dir = os.path.join(project_dir, "assets")
os.makedirs(assets_dir, exist_ok=True)

# Copy actual screenshots from gemini conversation directory if present
brain_dir = r"C:\Users\DELL\.gemini\antigravity-ide\brain\1671c805-0e4b-4d42-bd94-7ac6082ead66"
screenshot_mapping = {
    "ui_empty.png": "latest_ui_empty_1787324231499.png",
    "ui_sandbox.png": "latest_ui_results_1787324446876.png",
    "latest_ui_scan.png": "latest_ui_results_1787324446876.png",
    "diff_highlights.png": "latest_ui_diff_1787324684761.png",
    "ocr_scan.png": "latest_ui_single_ocr_1787325431264.png",
    "scanned_notes.jpg": ".user_uploaded/media_1786947808389.jpg",
    "system_workflow_diagram.jpg": "system_workflow_diagram_1787391164883.jpg",
    "college_logo.png": ".user_uploaded/media_1787391493749.png"
}

copied_screenshots = {}
for dest_name, src_name in screenshot_mapping.items():
    src_path = os.path.join(brain_dir, src_name)
    if os.path.exists(src_path):
        dest_path = os.path.join(assets_dir, dest_name)
        shutil.copy(src_path, dest_path)
        copied_screenshots[dest_name] = dest_path
        print(f"Copied screenshot: {dest_name}")
    else:
        # Fallback to local placeholders or look for similar files
        print(f"Screenshot {src_name} not found in brain directory.")

# --------------------------------------------------
# INITIALIZE PPTX PRESENTATION (16:9 Widescreen)
# --------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(15, 23, 42)      # Deep Navy (#0f172a)
CYAN = RGBColor(56, 189, 248)       # Neon Cyan (#38bdf8)
WHITE = RGBColor(255, 255, 255)
SILVER = RGBColor(148, 163, 184)    # Soft Gray (#94a3b8)
SLATE_CARD = RGBColor(30, 41, 59)   # Slate Card (#1e293b)
CODE_BG = RGBColor(10, 15, 29)      # Charcoal Code (#0a0f1d)

def set_slide_background(slide, color=DARK_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, color=CYAN, size=32):
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

def add_bullets(slide, left, top, width, height, bullets, size=18, color=SILVER):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    for idx, bullet in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.name = 'Segoe UI'
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(12)
        # Check if sub-bullet
        if bullet.strip().startswith("-") or bullet.strip().startswith("•"):
            p.level = 1
            p.font.size = Pt(size - 2)

def colorize_code_line(p, line, font_size):
    import re
    p.space_after = Pt(2)
    trimmed = line.strip()
    # Style comments completely in gray
    if trimmed.startswith("//") or trimmed.startswith("#"):
        run = p.add_run()
        run.text = line
        run.font.name = 'Consolas'
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(100, 116, 139)  # Slate Gray
        return

    # Regex to split line into strings, comments, numbers, keywords, and syntax tokens
    token_pattern = re.compile(r'(\".*?\"|\'.*?\'|\/\/.*|[\w]+|[^\w\s\x00-\x1f]+|\s+)')
    tokens = token_pattern.findall(line)
    
    keywords = {'async', 'function', 'try', 'catch', 'const', 'let', 'if', 'else', 'return', 'new', 'await', 'resolve', 'reject', 'then'}
    
    for token in tokens:
        run = p.add_run()
        run.text = token
        run.font.name = 'Consolas'
        run.font.size = Pt(font_size)
        
        # Color coding logic
        if token.startswith('"') or token.startswith("'"):
            run.font.color.rgb = RGBColor(253, 186, 116)  # Warm Orange for Strings
        elif token in keywords:
            run.font.color.rgb = RGBColor(244, 114, 182)  # Pink for Keywords
            run.font.bold = True
        elif token.startswith("//") or token.startswith("#"):
            run.font.color.rgb = RGBColor(100, 116, 139)  # Slate Gray for Comments
        elif token.isdigit():
            run.font.color.rgb = RGBColor(129, 140, 248)  # Indigo for Numbers
        elif token in {'document', 'window', 'Tesseract', 'pdfjsLib', 'mammoth', 'circle', 'prs'}:
            run.font.color.rgb = RGBColor(56, 189, 248)   # Cyan for built-in APIs
        else:
            run.font.color.rgb = RGBColor(241, 245, 249)  # Off-White for identifiers and punctuation

def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(71, 85, 105)
    shape.line.width = Pt(1)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)
    for idx, line in enumerate(code_lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        colorize_code_line(p, line, font_size)

def add_screenshot(slide, left, top, width, height, filename):
    img_path = os.path.join(assets_dir, filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width=width, height=height)
        print(f"Placed screenshot {filename} on slide.")
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = SLATE_CARD
        shape.line.color.rgb = RGBColor(71, 85, 105)
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = f"📷 [Actual Screenshot: {filename}]\nNot available in the current environment"
        p.font.name = 'Segoe UI'
        p.font.size = Pt(14)
        p.font.color.rgb = SILVER

# ==================================================
# SLIDE 1: TITLE SLIDE
# ==================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide1)

title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(0.2), Inches(11.83), Inches(1.4))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "AI-Based Document Similarity & Duplicate Detection System Using NLP"
p.font.name = 'Segoe UI'
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = CYAN

p_sub = tf.add_paragraph()
p_sub.alignment = PP_ALIGN.CENTER
p_sub.text = "AI-Based Document Processing Using OCR, NLP and Similarity Analysis\nPBL Task 2"
p_sub.font.name = 'Segoe UI'
p_sub.font.size = Pt(16)
p_sub.font.color.rgb = SILVER
p_sub.space_before = Pt(4)

# Subject Text Block
subj_box = slide1.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(0.8))
tf_subj = subj_box.text_frame
tf_subj.word_wrap = True
p_subj = tf_subj.paragraphs[0]
p_subj.alignment = PP_ALIGN.CENTER

run1 = p_subj.add_run()
run1.text = "Subject: "
run1.font.name = 'Segoe UI'
run1.font.size = Pt(18)
run1.font.bold = True
run1.font.color.rgb = RGBColor(163, 230, 53)  # Lime Yellow

run2 = p_subj.add_run()
run2.text = "Python for Data Science (BE05000231)"
run2.font.name = 'Segoe UI'
run2.font.size = Pt(18)
run2.font.bold = True
run2.font.color.rgb = WHITE

# Card block for student details
card = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.7), Inches(11.83), Inches(4.1))
card.fill.solid()
card.fill.fore_color.rgb = SLATE_CARD
card.line.color.rgb = RGBColor(71, 85, 105)
card.line.width = Pt(1.5)

# Student 1
s1_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.9), Inches(5.0), Inches(1.1))
tf_s1 = s1_box.text_frame
tf_s1.word_wrap = True
p_s1_n = tf_s1.paragraphs[0]
r1_n_k = p_s1_n.add_run()
r1_n_k.text = "NAME 1: "
r1_n_k.font.bold = True
r1_n_k.font.color.rgb = CYAN
r1_n_v = p_s1_n.add_run()
r1_n_v.text = "Nimavat Jaydeep M."
r1_n_v.font.color.rgb = WHITE

p_s1_e = tf_s1.add_paragraph()
p_s1_e.space_before = Pt(4)
r1_e_k = p_s1_e.add_run()
r1_e_k.text = "ENROLLMENT 1: "
r1_e_k.font.bold = True
r1_e_k.font.color.rgb = CYAN
r1_e_v = p_s1_e.add_run()
r1_e_v.text = "250043107032"
r1_e_v.font.color.rgb = WHITE

# Student 2
s2_box = slide1.shapes.add_textbox(Inches(6.8), Inches(2.9), Inches(5.0), Inches(1.1))
tf_s2 = s2_box.text_frame
tf_s2.word_wrap = True
p_s2_n = tf_s2.paragraphs[0]
r2_n_k = p_s2_n.add_run()
r2_n_k.text = "NAME 2: "
r2_n_k.font.bold = True
r2_n_k.font.color.rgb = CYAN
r2_n_v = p_s2_n.add_run()
r2_n_v.text = "Upadhyay Hitarth S."
r2_n_v.font.color.rgb = WHITE

p_s2_e = tf_s2.add_paragraph()
p_s2_e.space_before = Pt(4)
r2_e_k = p_s2_e.add_run()
r2_e_k.text = "ENROLLMENT 2: "
r2_e_k.font.bold = True
r2_e_k.font.color.rgb = CYAN
r2_e_v = p_s2_e.add_run()
r2_e_v.text = "250043107050"
r2_e_v.font.color.rgb = WHITE

for r in [r1_n_k, r1_n_v, r1_e_k, r1_e_v, r2_n_k, r2_n_v, r2_e_k, r2_e_v]:
    r.font.name = 'Segoe UI'
    r.font.size = Pt(17)

# Class Details
class_box = slide1.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(10.9), Inches(0.6))
tf_class = class_box.text_frame
p_class = tf_class.paragraphs[0]
p_class.alignment = PP_ALIGN.CENTER
r_br_k = p_class.add_run()
r_br_k.text = "BRANCH: "
r_br_k.font.bold = True
r_br_k.font.color.rgb = CYAN
r_br_v = p_class.add_run()
r_br_v.text = "Computer Engineering     |     "
r_br_v.font.color.rgb = WHITE

r_sem_k = p_class.add_run()
r_sem_k.text = "SEMESTER: "
r_sem_k.font.bold = True
r_sem_k.font.color.rgb = CYAN
r_sem_v = p_class.add_run()
r_sem_v.text = "5th"
r_sem_v.font.color.rgb = WHITE

for r in [r_br_k, r_br_v, r_sem_k, r_sem_v]:
    r.font.name = 'Segoe UI'
    r.font.size = Pt(17)

# College Logo Container
logo_bg = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.66), Inches(5.1), Inches(6.0), Inches(1.3))
logo_bg.fill.solid()
logo_bg.fill.fore_color.rgb = WHITE
logo_bg.line.color.rgb = CYAN
logo_bg.line.width = Pt(1.5)

logo_path = os.path.join(assets_dir, "college_logo.png")
if os.path.exists(logo_path):
    slide1.shapes.add_picture(logo_path, Inches(3.76), Inches(5.15), width=Inches(5.8), height=Inches(1.2))
else:
    # Text Fallback if file missing
    coll_box = slide1.shapes.add_textbox(Inches(1.2), Inches(5.1), Inches(10.9), Inches(0.6))
    tf_coll = coll_box.text_frame
    p_coll = tf_coll.paragraphs[0]
    p_coll.alignment = PP_ALIGN.CENTER
    r_coll_k = p_coll.add_run()
    r_coll_k.text = "COLLEGE: "
    r_coll_k.font.bold = True
    r_coll_k.font.color.rgb = CYAN
    r_coll_v = p_coll.add_run()
    r_coll_v.text = "B H Gardi College Of Engineering And Technology"
    r_coll_v.font.color.rgb = RGBColor(15, 23, 42) # dark text on white fallback container if drawn
    for r in [r_coll_k, r_coll_v]:
        r.font.name = 'Segoe UI'
        r.font.size = Pt(17)

# ==================================================
# SLIDE 2: INTRODUCTION
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "1. Introduction")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Document Similarity: The process of measuring the semantic or syntactic alignment between two text bodies.",
    "• Duplicate Detection: Identifying exact or highly rewritten copies of files to protect academic integrity.",
    "• Natural Language Processing (NLP): Used to clean, tokenize, and map raw text characters into mathematical arrays.",
    "• Automated Comparison: Bypasses slow manual grading reviews by running comparison algorithms in real-time.",
    "• Scope of Present Work: A client-side, browser-native pipeline capable of parsing PDF, DOCX, and scanned note photos."
])

# ==================================================
# SLIDE 3: PROJECT OVERVIEW
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "2. Project Overview")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(6.0), Inches(4.8), [
    "• 100% Client-Side Web Application: Executed entirely in the browser using HTML5, CSS3, and JavaScript.",
    "• Document Ingestion: Extracts text dynamically from pasted strings or uploaded files (.txt, .pdf, .docx).",
    "• Offline OCR Scanner: Extracts handwritten or printed characters from note photos (.png, .jpg, .jpeg).",
    "• Statistical NLP Math: Matches vocabulary sets using Cosine Vector models and Jaccard Set indices.",
    "• Immediate Feedback: Shows circular gauges, metrics, and Turnitin-style highlighted match segments."
])
add_screenshot(slide, Inches(7.2), Inches(1.8), Inches(5.38), Inches(4.2), "ui_empty.png")

# ==================================================
# SLIDE 4: PROBLEM STATEMENT
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "3. Problem Statement")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Inability to Scan Images: Standard plagiarism detectors completely fail to inspect book photos or handwritten notebook scans.",
    "• Network Dependency Constraints: Python server setups are blocked by university firewall sockets, causing terminal timeout errors.",
    "• Heavy Server Costs: Spawning GPU-dependent server-side OCR modules is expensive and complex to maintain.",
    "• Data Leaks and Exposure: Uploading student work to proprietary databases exposes sensitive student files."
])

# ==================================================
# SLIDE 5: PROJECT OBJECTIVES
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "4. Project Objectives")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Serverless Architecture: Deliver a high-performance scanner running completely on local client browser threads.",
    "• Client-Side OCR Scanning: Perform text extractions from student note photos locally using Web Workers.",
    "• Combined Mathematical Models: Implement TF-IDF Cosine Similarity and Jaccard Indexing to catch rewritten copies.",
    "• Turnitin-style Diff: Render real-time visual highlight tags showing exact matching vocabularies.",
    "• Maximum Data Privacy: Prevent documents from ever leaving the student's browser."
])

# ==================================================
# SLIDE 6: PROJECT MOTIVATION
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "5. Project Motivation")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Zero Maintenance Overhead: Developing static pages (Netlify-ready) that run forever without paid database backends.",
    "• Bypassing Proxy Blocks: Eliminating Python terminal connection errors (`WinError 10060`) by using browser WebAssembly.",
    "• Time-Saving Automation: Evaluating handwritten pages or mathematical outlines in seconds.",
    "• Mobile Accessibility: Creating responsive layouts that fit on tablets and phones in a presentation room."
])

# ==================================================
# SLIDE 7: EXISTING SYSTEM / EXISTING APPROACH
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "6. Existing System / Existing Approach")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Python Flask/Streamlit Backends: Spawns heavy local servers running PyTorch or OpenCV frameworks.",
    "• Paid Subscriptions: Academic institutions depend on proprietary APIs (Copyleaks, Turnitin) with strict rate limits.",
    "• Static String Checks: Basic checkers match exact character lists, ignoring vocabulary distributions or synonym rewriting.",
    "• Manual Student Uploads: Teachers have to manually copy and paste image text before running checks."
])

# ==================================================
# SLIDE 8: LIMITATIONS OF EXISTING SYSTEM
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "7. Limitations of Existing System")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Network Timout Errors: Server socket connections crash frequently during model weight downloads.",
    "• Latency: Transmitting high-resolution student images to external APIs causes long processing queues.",
    "• Privacy Non-Compliance: Storing documents in shared clouds violates university student privacy policies.",
    "• Lack of Offline Scans: stream processes fail completely if the server loses internet connection."
])

# ==================================================
# SLIDE 9: PROPOSED SYSTEM
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "8. Proposed System")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(6.0), Inches(4.8), [
    "• 100% Client-Side Engine: Browser WebAssembly compiles local parser scripts.",
    "• Multi-Format File Uploads: PDF.js and Mammoth.js process files locally.",
    "• IndexedDB Model Caching: Stores Tesseract OCR neural weights inside browser cache.",
    "• JavaScript NLP execution: tokenizes, cleans, and runs vector distance calculations on-device."
])
add_screenshot(slide, Inches(7.2), Inches(1.8), Inches(5.38), Inches(4.2), "ui_sandbox.png")

# ==================================================
# SLIDE 10: KEY FEATURES
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "9. Key Features")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• WebAssembly Tesseract OCR: Scans JPG/PNG text directly in the browser.",
    "• PDF & Word Extraction: Parses files page-by-page client-side using JavaScript array buffers.",
    "• Animated SVG Circular Gauge: Dynamic progress dial tracking the similarity score.",
    "• Sandbox Demo presets: Pre-loaded group theory notes to test features in one click.",
    "• Citation Filter checkbox: Ignores text inside double/single quotes during scoring."
])

# ==================================================
# SLIDE 11: PROJECT SCOPE
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "10. Project Scope")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Academic Grading Portals: Fast assignment checks directly within learning management portals.",
    "• Offline Verification: Scanning physical textbook sheets in offline university labs.",
    "• Private Homework Submissions: Validates original student answers without exposing them online.",
    "• Lightweight Plagiarism Review: Useful for grading committees to check duplicate reports on their phones."
])

# ==================================================
# SLIDE 12: TECHNOLOGY STACK
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "11. Technology Stack")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Core Technologies: HTML5 Semantic Layout, CSS3 Variables (theme layout), ES6 JavaScript.",
    "• OCR Parser: Tesseract.js (v5, runs Web Workers on local browser threads).",
    "• Adobe PDF Parser: PDF.js (v3.11.174, page-coordinate array parser).",
    "• Word Document Parser: Mammoth.js (v1.6.0, docx-to-XML plain text converter).",
    "• Hosting Platform: Netlify Continuous Deployment (connected via GitHub)."
])

# ==================================================
# SLIDE 13: SYSTEM ARCHITECTURE
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "12. System Architecture")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• User Interface: Glassmorphic HTML5 dashboard handles inputs (files/text).",
    "• Extraction Layer: FileReader API reads bytes locally. Tesseract, PDF.js, and Mammoth parse to strings.",
    "• NLP Core: JavaScript tokenizes text, cleans punctuation, and normalizes cases.",
    "• Comparison Core: JavaScript vectors compute Cosine and Jaccard similarity indices.",
    "• Rendering: Updates orange highlighter spans and circular progress offsets dynamically."
])

# ==================================================
# SLIDE 14: SYSTEM WORKFLOW
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "13. System Workflow")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(6.0), Inches(4.8), [
    "• Select: User selects text files, PDF files, Word files, or image scans.",
    "• Loading: Spawns progress spinner overlays on form submit.",
    "• Extract: Browser libraries parse text content on workers.",
    "• Calculate: Computes Cosine + Jaccard average match.",
    "• Render: Updates UI gauge and diff spans instantly."
])
add_screenshot(slide, Inches(7.2), Inches(1.8), Inches(5.38), Inches(4.2), "system_workflow_diagram.jpg")

# ==================================================
# SLIDE 15: INPUT DATA / DOCUMENT INPUT
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "14. Input Data / Document Input")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Uploading File Formats: Supports .txt, .pdf, .docx, and images (.png, .jpg, .jpeg).",
    "• Base64 Image Preview: Shows thumbnail previews of uploaded images directly in the columns.",
    "• Text Box Copy-Pasting: Users can type directly into the text areas to run instant comparisons.",
    "• Sandbox Note Presets: Clickable buttons to load university math notes for quick validation."
])

# ==================================================
# SLIDE 16: DOCUMENT PROCESSING
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "15. Document Processing")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• FileReader API Conversion: Reads files locally using Javascript array buffers.",
    "• PDF.js Document parsing: `pdfjsLib.getDocument()` loads PDF arrays to memory.",
    "• Mammoth.js Paragraph extraction: Reads .docx file formats directly into clean text strings.",
    "• Web Worker sandboxing: Parsing processes run in browser threads so the dashboard stays fast."
])

# ==================================================
# SLIDE 17: TEXT EXTRACTION / PREPROCESSING
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "16. Text Preprocessing Pipeline")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Case Normalization: `.toLowerCase()` standardizes casing to prevent capitalization mismatches.",
    "• Punctuation Removal: RegEx `.replace(/[^\w\s]/g, '')` strips symbols, commas, and periods.",
    "• Space Trimming: `.trim()` removes leading and trailing spaces.",
    "• Word Splitting: `.split(/\\s+/)` separates the cleaned strings into token arrays of clean words."
])

# ==================================================
# SLIDE 18: NLP PROCESSING
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "17. NLP Processing")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Tokenization: Splits text strings into arrays of individual word token strings.",
    "• Vocabulary Set indexing: Converts token arrays into unique JavaScript Set structures.",
    "• Frequency mapping: Counts word occurrences to build frequency bags for vector matching.",
    "• Citation Filtering: Regex replaces quoted sections with blank space before similarity checks."
])

# ==================================================
# SLIDE 19: TEXT REPRESENTATION
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "18. Text Representation")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Vocabulary Union Vector: Compares word sets from both documents to create a master list of all unique words.",
    "• Term Frequency (TF) Maps: Maps occurrences of each word coordinate relative to the master list.",
    "• Multi-Dimensional Vectors: Represents both documents as coordinates, ready for distance calculation."
])

# ==================================================
# SLIDE 20: TF-IDF IMPLEMENTATION
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "19. Term Frequency (TF) Vector Method")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Actual Method in Code: Dynamic Word Frequency Vectors.",
    "• Vocab Set: Combines both documents' word arrays into a master Vocab Set.",
    "• Frequency Mapping: Loops vocab list to count term occurrences in each document.",
    "• Vector Representation: Maps counts into coordinate arrays, ready for Cosine Distance calculations."
])

# ==================================================
# SLIDE 21: COSINE SIMILARITY
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "20. Cosine Similarity Vector Math")

add_bullets(slide, Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8), [
    "• Cosine Similarity measures the angle between document vectors:",
    "  Cosine Similarity = (A · B) / (||A|| * ||B||)",
    "• Dot Product: Multiplies matching coordinates between document maps.",
    "• Magnitudes product: Normalizes vectors to keep calculations accurate.",
    "• Length Independence: Ensures matching remains accurate even if document lengths differ."
])

code_cosine = [
    "// Vector math in calculateCosine()",
    "let dotProduct = 0;",
    "let mag1 = 0;",
    "let mag2 = 0;",
    "vocab.forEach(word => {",
    "    dotProduct += freq1[word] * freq2[word];",
    "    mag1 += freq1[word] * freq1[word];",
    "    mag2 += freq2[word] * freq2[word];",
    "});",
    "if (mag1 === 0 || mag2 === 0) return 0;",
    "return (dotProduct / (Math.sqrt(mag1) * Math.sqrt(mag2))) * 100;"
]
add_code_block(slide, Inches(6.5), Inches(1.8), Inches(6.08), Inches(4.2), code_cosine, font_size=11)

# ==================================================
# SLIDE 22: SIMILARITY DETECTION
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "21. Similarity Detection (Jaccard Index)")

add_bullets(slide, Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8), [
    "• Checks overall Vocabulary Set overlap:",
    "  Jaccard Index = |A ∩ B| / |A ∪ B|",
    "• Intersection: Counts shared unique words.",
    "• Union: Counts total combined unique words.",
    "• Dividing intersection by union outputs Jaccard percentage.",
    "• Highly effective at catching copied vocabulary sets."
])

code_jaccard = [
    "function calculateJaccard(text1, text2) {",
    "    const t1 = cleanText(text1);",
    "    const t2 = cleanText(text2);",
    "    if (!t1 || !t2) return 0;",
    "    const words1 = new Set(t1.split(/\\s+/));",
    "    const words2 = new Set(t2.split(/\\s+/));",
    "    const intersection = new Set(",
    "        [...words1].filter(x => words2.has(x))",
    "    );",
    "    const union = new Set([...words1, ...words2]);",
    "    if (union.size === 0) return 0;",
    "    return (intersection.size / union.size) * 100;",
    "}"
]
add_code_block(slide, Inches(6.5), Inches(1.8), Inches(6.08), Inches(4.2), code_jaccard, font_size=10.5)

# ==================================================
# SLIDE 23: DUPLICATE DETECTION LOGIC
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "22. Duplicate Detection Logic & Thresholds")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Combined Score: Calculates the average of Cosine Similarity and Jaccard Index score.",
    "• Severity Warning Classifications in JavaScript:",
    "  - 90% - 100%: Direct Copy / Duplicate (Red Color)",
    "  - 70% - 89%: Highly Similar (Orange Color)",
    "  - 40% - 69%: Moderately Similar (Yellow Color)",
    "  - 15% - 39%: Slightly Similar (Blue Color)",
    "  - Below 15%: Completely Unique (Green Color)"
])

# ==================================================
# SLIDE 24: DOCUMENT COMPARISON
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "23. Document Comparison Preview")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(6.0), Inches(4.8), [
    "• Visual Word Diff: Highlight algorithm checks if a word in Document A exists in the vocabulary set of Document B.",
    "• Matches are wrapped in custom span elements.",
    "• Orange Highlights: Identifies identical matching words in both documents.",
    "• Clipboard integration: Copy buttons to extract matching strings in one click."
])
add_screenshot(slide, Inches(7.2), Inches(1.8), Inches(5.38), Inches(4.2), "diff_highlights.png")

# ==================================================
# SLIDE 25: SIMILARITY RESULTS
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "24. Similarity Results Dashboard")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(6.0), Inches(4.8), [
    "• Custom Score Gauge: Shows final similarity score inside the dynamic progress circle.",
    "• Match Status Card: Displays text descriptions (e.g. 'Moderately Similar', 'Highly Similar').",
    "• Word Count Stats: Renders exact word counts for both Document A and Document B.",
    "• Visual preview: Displays the uploaded note photo directly in the column."
])
add_screenshot(slide, Inches(7.2), Inches(1.8), Inches(5.38), Inches(4.2), "latest_ui_scan.png")

# ==================================================
# SLIDE 26: SIMILARITY VISUALIZATION
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "25. Similarity Visualization (SVG Gauge)")

add_bullets(slide, Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8), [
    "• Circular progress ring rendered in inline SVG.",
    "• Circumference = 326.7px (radius = 52px).",
    "• JS computes dynamic stroke-dashoffset for score percentage.",
    "• Smooth 1.5s CSS transition filling animation.",
    "• Glow and circle colors update dynamically matching similarity severity levels."
])

code_gauge = [
    "// Render score with SVG gauge animation",
    "const circle = document.getElementById('radial-bar');",
    "const circumference = 326.7;",
    "const offset = circumference - (circumference * scoreInt / 100);",
    "circle.style.strokeDashoffset = offset;",
    "circle.style.stroke = color;",
    "circle.style.filter = `drop-shadow(0 0 8px ${color}80)`;"
]
add_code_block(slide, Inches(6.5), Inches(1.8), Inches(6.08), Inches(4.2), code_gauge, font_size=12)

# ==================================================
# SLIDE 27: USER INTERFACE
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "26. User Interface Design")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(6.0), Inches(4.8), [
    "• Midnight Slate-Blue gradient background (`#0f172a` to `#1e293b`).",
    "• Glassmorphic interface: Semi-transparent panels with background blur.",
    "• Google Fonts integration (Plus Jakarta Sans).",
    "• Dual Column Previews with uploaded image display placeholders.",
    "• Single OCR scanner mode and double comparison mode."
])
add_screenshot(slide, Inches(7.2), Inches(1.8), Inches(5.38), Inches(4.2), "latest_ui_scan.png")

# ==================================================
# SLIDE 28: IMPLEMENTATION / CODE HIGHLIGHTS
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "27. Main Web Controller Logic")

add_bullets(slide, Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8), [
    "• Client-side controller routing in JS.",
    "• Handles extraction, OCR, comparisons.",
    "• Single scanner vs Dual match flow check.",
    "• Shows loader overlay during calculations.",
    "• Catches errors and updates error container."
])

code_main = [
    "async function startScanner() {",
    "    const overlay = document.getElementById('loader-overlay');",
    "    const errorBox = document.getElementById('error-box');",
    "    errorBox.style.display = 'none';",
    "    overlay.style.display = 'flex';",
    "    try {",
    "        const t1 = await extractText('file1', 'text1', 'Document A');",
    "        const t2 = await extractText('file2', 'text2', 'Document B');",
    "        if (t1.trim() && !t2.trim()) {",
    "            renderSingleScannerResult('A', t1, file1);",
    "        } else {",
    "            renderComparisonResult(t1, t2);",
    "        }",
    "    } catch (err) {",
    "        errorBox.innerText = err;",
    "    }",
    "}"
]
add_code_block(slide, Inches(6.5), Inches(1.8), Inches(6.08), Inches(4.2), code_main, font_size=10)

# ==================================================
# SLIDE 29: TESTING & VALIDATION
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "28. Testing & Validation")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(6.0), Inches(4.8), [
    "• Tested offline OCR extraction with handwritten note scans (Math mappings and functions).",
    "• Tested digital PDF extraction page-by-page to check text stream reconstruction.",
    "• Verified similarity scores using identical and modified notes in the sandbox selector.",
    "• Confirmed that network firewalls have zero impact on scanning speeds or execution."
])
add_screenshot(slide, Inches(7.2), Inches(1.8), Inches(5.38), Inches(4.2), "scanned_notes.jpg")

# ==================================================
# SLIDE 30: ADVANTAGES & LIMITATIONS
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "29. Advantages & Limitations")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Advantages:",
    "  - 100% Serverless: Bypasses server maintenance fees and offline/online proxy socket blocks.",
    "  - Data Privacy: Files are processed completely in the browser and never sent online.",
    "  - Mobile Responsive: Renders cleanly on phones and tablets during presentations.",
    "• Limitations:",
    "  - CPU Bound: Heavy files take longer to scan via client-side JavaScript threads.",
    "  - No Database: Similarity reports cannot be saved in a remote history database."
])

# ==================================================
# SLIDE 31: FUTURE SCOPE
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)
add_title(slide, "30. Future Scope")
add_bullets(slide, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), [
    "• Multi-Language OCR: Add dropdown support for Spanish, Hindi, or French training models.",
    "• Web Link comparisons: Paste direct URLs to fetch and check similarity against live websites.",
    "• PDF report downloads: Compile plagiarism match reports directly in the browser to download offline.",
    "• Sentence Embeddings: Integrate client-side Transformers (like ONNX) for semantic similarity matching."
])

# ==================================================
# SLIDE 32: CONCLUSION
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)

# Header
title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "31. Conclusion"
p.font.name = 'Segoe UI'
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = CYAN

# Body
body_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.83), Inches(2.2))
tf_body = body_box.text_frame
tf_body.word_wrap = True
p_b = tf_body.paragraphs[0]
p_b.text = "We have successfully built a fully serverless, client-side document similarity and duplicate detection system using Tesseract.js, PDF.js, and Mammoth.js.\n\nThe project demonstrates that advanced AI parsing and NLP comparison algorithms can run completely locally on consumer browser engines, eliminating the need for expensive GPU-hosting setups."
p_b.font.name = 'Segoe UI'
p_b.font.size = Pt(20)
p_b.font.color.rgb = SILVER
p_b.space_after = Pt(15)

# Thank You
thanks_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.8), Inches(11.83), Inches(1.8))
tf_t = thanks_box.text_frame
p_t = tf_t.paragraphs[0]
p_t.alignment = PP_ALIGN.CENTER
p_t.text = "THANK YOU\n\nQuestions?"
p_t.font.name = 'Segoe UI'
p_t.font.size = Pt(36)
p_t.font.bold = True
p_t.font.color.rgb = CYAN

# --------------------------------------------------
# SAVE PRESENTATION (To both output file paths with robust fallbacks)
# --------------------------------------------------
import time
dest_filenames = [
    "AI_Based_Document_Similarity_Duplicate_Detection_NLP_PBL",
    "AI_Document_OCR_Scanner_Similarity_Detector_PBL"
]

for base_name in dest_filenames:
    out_path = os.path.join(project_dir, f"{base_name}.pptx")
    saved = False
    try:
        prs.save(out_path)
        print(f"Presentation saved successfully to {out_path}")
        saved = True
    except PermissionError:
        pass

    if not saved:
        alt_path = os.path.join(project_dir, f"{base_name}_v2.pptx")
        try:
            prs.save(alt_path)
            print(f"WARNING: original file locked. Saved alternative version successfully to: {alt_path}")
            saved = True
        except PermissionError:
            pass

    if not saved:
        timestamp = int(time.time())
        timestamp_path = os.path.join(project_dir, f"{base_name}_{timestamp}.pptx")
        prs.save(timestamp_path)
        print(f"WARNING: Both original and v2 files are locked! Saved with timestamp to: {timestamp_path}")
